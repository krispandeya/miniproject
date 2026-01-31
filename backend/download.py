from docx import Document
from pptx import Presentation
import pandas as pd
import csv
import io
from fpdf import FPDF

from text import extract_text_from_image


def build_export(text, file_type):
    file_type = (file_type or "pdf").lower()

    if file_type == "pdf":
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)
        for line in text.splitlines() or [text]:
            pdf.multi_cell(0, 8, line)
        data = pdf.output(dest="S").encode("latin-1")
        return "scan.pdf", "application/pdf", data

    if file_type == "docx":
        buf = io.BytesIO()
        doc = Document()
        doc.add_paragraph(text)
        doc.save(buf)
        return "scan.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", buf.getvalue()

    if file_type == "pptx":
        buf = io.BytesIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.placeholders[1].text = text
        prs.save(buf)
        return "scan.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", buf.getvalue()

    if file_type == "xlsx":
        buf = io.BytesIO()
        pd.DataFrame([{"Content": text}]).to_excel(buf, index=False)
        return "scan.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", buf.getvalue()

    if file_type == "csv":
        buf = io.StringIO()
        csv.writer(buf).writerow([text])
        data = buf.getvalue().encode("utf-8")
        return "scan.csv", "text/csv", data

    if file_type == "html":
        data = f"<html><body><p>{text}</p></body></html>".encode("utf-8")
        return "scan.html", "text/html", data

    if file_type == "txt":
        data = text.encode("utf-8")
        return "scan.txt", "text/plain", data

    return None, None, None


if __name__ == "__main__":
    user_choice = input("Enter format (docx, pptx, pdf, csv, html, txt, or xlsx): ").lower()
    text = extract_text_from_image()
    filename, _, data = build_export(text, user_choice)
    if not data:
        print("Invalid choice! Please run again and pick docx, pptx, pdf, xlsx, csv, html, or txt.")
    else:
        with open(filename, "wb") as f:
            f.write(data)
        print(f"Finished! Your {user_choice} file is ready.")